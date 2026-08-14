from __future__ import annotations

import json
import mimetypes
import re
import shutil
import subprocess
import tarfile
import tempfile
import textwrap
import zipfile
from pathlib import Path
from typing import Any

import pandas as pd
import yaml
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4, portrait
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from pypdf import PdfReader
import fitz  # PyMuPDF
from xml.etree import ElementTree as ET

from utils.logger import logger
from utils.temp_files import (
    stem_of,
    make_output_name,
)

# -----------------------------
# Feature detection
# -----------------------------
HAS_LIBREOFFICE = shutil.which("libreoffice") or shutil.which("soffice")
HAS_FFMPEG = shutil.which("ffmpeg")
HAS_INKSCAPE = shutil.which("inkscape")
HAS_MAGICK = shutil.which("magick") or shutil.which("convert")

# -----------------------------
# Categories
# -----------------------------
IMAGE_EXTS = {"png", "jpg", "jpeg", "webp", "gif", "bmp", "tif", "tiff", "svg"}
TEXT_EXTS = {"txt", "json", "csv", "xml", "yaml", "yml"}
SHEET_EXTS = {"xlsx", "ods"}
DOC_EXTS = {"pdf", "docx", "odt", "rtf", "doc"}
PRESENTATION_EXTS = {"pptx", "ppt", "odp"}
VIDEO_EXTS = {"mp4", "mov", "avi", "mkv", "webm"}
AUDIO_EXTS = {"mp3", "wav", "ogg", "flac"}
ARCHIVE_EXTS = {"zip"}

IMAGE_TARGETS = {"png", "jpg", "jpeg", "webp", "gif", "bmp", "tif", "pdf", "docx", "pptx"}
TEXT_TARGETS = {"txt", "json", "csv", "xml", "yaml", "xlsx", "ods", "pdf", "docx", "pptx", "png"}
SHEET_TARGETS = {"xlsx", "ods", "csv", "json", "txt", "pdf", "docx", "pptx", "png"}
DOC_TARGETS = {"pdf", "docx", "txt", "csv", "json", "png", "xlsx", "pptx"}
PPT_TARGETS = {"pptx", "pdf", "txt", "csv", "json", "png", "docx", "xlsx"}
VIDEO_TARGETS = {"mp4", "mov", "avi", "mkv", "webm", "mp3", "wav", "ogg", "flac"}
AUDIO_TARGETS = {"mp3", "wav", "ogg", "flac", "mp4", "webm"}
ARCHIVE_TARGETS = {"extract"}

SUPPORTED_CONVERSIONS: dict[str, list[str]] = {
    "png": sorted(IMAGE_TARGETS),
    "jpg": sorted(IMAGE_TARGETS),
    "jpeg": sorted(IMAGE_TARGETS),
    "webp": sorted(IMAGE_TARGETS),
    "gif": sorted(IMAGE_TARGETS),
    "bmp": sorted(IMAGE_TARGETS),
    "svg": ["png", "jpg", "jpeg", "webp", "gif", "bmp", "tif", "pdf", "docx", "pptx"],
    "tif": sorted(IMAGE_TARGETS),
    "tiff": sorted(IMAGE_TARGETS),

    "txt": sorted(TEXT_TARGETS),
    "json": sorted(TEXT_TARGETS),
    "csv": sorted(TEXT_TARGETS),
    "xml": sorted(TEXT_TARGETS),
    "yaml": sorted(TEXT_TARGETS),
    "yml": sorted(TEXT_TARGETS),

    "xlsx": sorted(SHEET_TARGETS),
    "ods": sorted(SHEET_TARGETS),

    "pdf": sorted(DOC_TARGETS),
    "docx": sorted(DOC_TARGETS),
    "odt": sorted(DOC_TARGETS),
    "rtf": sorted(DOC_TARGETS),
    "doc": sorted(DOC_TARGETS),

    "pptx": sorted(PPT_TARGETS),
    "ppt": sorted(PPT_TARGETS),
    "odp": sorted(PPT_TARGETS),

    "mp4": sorted(VIDEO_TARGETS),
    "mov": sorted(VIDEO_TARGETS),
    "avi": sorted(VIDEO_TARGETS),
    "mkv": sorted(VIDEO_TARGETS),
    "webm": sorted(VIDEO_TARGETS),

    "mp3": sorted(AUDIO_TARGETS),
    "wav": sorted(AUDIO_TARGETS),
    "ogg": sorted(AUDIO_TARGETS),
    "flac": sorted(AUDIO_TARGETS),

    "zip": sorted(ARCHIVE_TARGETS),
}

MIME_BY_EXT = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "webp": "image/webp",
    "gif": "image/gif",
    "bmp": "image/bmp",
    "tif": "image/tiff",
    "tiff": "image/tiff",
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "ods": "application/vnd.oasis.opendocument.spreadsheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "csv": "text/csv",
    "json": "application/json",
    "xml": "application/xml",
    "yaml": "text/yaml",
    "yml": "text/yaml",
    "txt": "text/plain",
    "mp3": "audio/mpeg",
    "wav": "audio/wav",
    "ogg": "audio/ogg",
    "flac": "audio/flac",
    "mp4": "video/mp4",
    "mov": "video/quicktime",
    "avi": "video/x-msvideo",
    "mkv": "video/x-matroska",
    "webm": "video/webm",
    "zip": "application/zip",
}

SUGGESTED_SHEETS = ["Sheet1"]
DEFAULT_PDF_PAGESIZE = A4

def source_category(ext: str) -> str:
    ext = (ext or "").lower().strip()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in TEXT_EXTS:
        return "text-data"
    if ext in SHEET_EXTS:
        return "spreadsheet"
    if ext in DOC_EXTS:
        return "document"
    if ext in PRESENTATION_EXTS:
        return "presentation"
    if ext in VIDEO_EXTS:
        return "video"
    if ext in AUDIO_EXTS:
        return "audio"
    if ext in ARCHIVE_EXTS:
        return "archive"
    return "unknown"

def normalise_ext(ext: str) -> str:
    ext = (ext or "").lower().strip().lstrip(".")
    return "jpg" if ext == "jpeg" else ("tif" if ext == "tiff" else ("yaml" if ext == "yml" else ext))

def has_feature(name: str) -> bool:
    return bool({"libreoffice": HAS_LIBREOFFICE, "ffmpeg": HAS_FFMPEG, "inkscape": HAS_INKSCAPE, "magick": HAS_MAGICK}.get(name))

def ensure_parent(path: Path):
    path.parent.mkdir(parents=True, exist_ok=True)

def safe_utf8_text(path: Path) -> str:
    data = path.read_bytes()
    for enc in ("utf-8-sig", "utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except Exception:
            continue
    return data.decode("utf-8", errors="replace")

def image_from_svg(svg_path: Path, out_png: Path) -> Path:
    ensure_parent(out_png)
    if HAS_INKSCAPE:
        cmd = ["inkscape", str(svg_path), "--export-type=png", f"--export-filename={out_png}"]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return out_png
    if HAS_MAGICK:
        cmd = ["magick", "convert", str(svg_path), str(out_png)] if Path(str(HAS_MAGICK)).name == "magick" else ["convert", str(svg_path), str(out_png)]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return out_png
    raise RuntimeError("SVG conversion requires Inkscape or ImageMagick.")

def load_image_any(path: Path) -> Image.Image:
    ext = normalise_ext(path.suffix)
    if ext == "svg":
        tmp_png = path.with_suffix(".render.png")
        image_from_svg(path, tmp_png)
        img = Image.open(tmp_png)
        return img.copy()
    img = Image.open(path)
    try:
        img.load()
    except Exception:
        pass
    if getattr(img, "is_animated", False):
        try:
            img.seek(0)
        except Exception:
            pass
    return img

def image_to_pdf(img: Image.Image, out_pdf: Path) -> Path:
    ensure_parent(out_pdf)
    if img.mode in ("RGBA", "P"):
        img = img.convert("RGB")
    img.save(out_pdf, "PDF", resolution=144.0)
    return out_pdf

def image_to_raster(img: Image.Image, out_path: Path, ext: str) -> Path:
    ensure_parent(out_path)
    ext = normalise_ext(ext)
    if img.mode == "P":
        img = img.convert("RGBA")
    if ext in {"jpg", "jpeg"}:
        img = img.convert("RGB")
        img.save(out_path, quality=92, optimize=True)
    elif ext == "png":
        img.save(out_path, optimize=True)
    elif ext == "webp":
        img.save(out_path, quality=90, method=6)
    elif ext == "gif":
        img = img.convert("P", palette=Image.Palette.ADAPTIVE)
        img.save(out_path)
    elif ext == "bmp":
        img.save(out_path)
    elif ext in {"tif", "tiff"}:
        img.save(out_path, compression="tiff_deflate")
    else:
        raise RuntimeError(f"Unsupported image output format: {ext}")
    return out_path

def fit_image_to_page(img: Image.Image, page_w: int, page_h: int, margin: int = 36) -> tuple[int, int, int, int]:
    w, h = img.size
    avail_w = page_w - margin * 2
    avail_h = page_h - margin * 2
    ratio = min(avail_w / w, avail_h / h)
    nw, nh = int(w * ratio), int(h * ratio)
    x = int((page_w - nw) / 2)
    y = int((page_h - nh) / 2)
    return x, y, nw, nh

def render_pdf_to_images(pdf_path: Path, out_dir: Path, base_name: str) -> list[Path]:
    doc = fitz.open(pdf_path)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_paths: list[Path] = []
    for idx, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        out = out_dir / f"{base_name}_page_{idx}.png"
        pix.save(out)
        out_paths.append(out)
    return out_paths

def zip_paths(paths: list[Path], out_zip: Path) -> Path:
    ensure_parent(out_zip)
    with zipfile.ZipFile(out_zip, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for p in paths:
            zf.write(p, arcname=p.name)
    return out_zip

def dataframe_from_any(data: Any) -> pd.DataFrame:
    if isinstance(data, pd.DataFrame):
        return data
    if isinstance(data, list):
        if not data:
            return pd.DataFrame({"value": []})
        if all(isinstance(item, dict) for item in data):
            return pd.DataFrame(data)
        return pd.DataFrame({"value": data})
    if isinstance(data, dict):
        # try to preserve dict structure
        if all(not isinstance(v, (dict, list)) for v in data.values()):
            return pd.DataFrame([data])
        rows = []
        for key, value in data.items():
            rows.append({"key": key, "value": json.dumps(value, ensure_ascii=False)})
        return pd.DataFrame(rows)
    if data is None:
        return pd.DataFrame({"value": []})
    return pd.DataFrame({"value": [str(data)]})

def parse_json_to_data(text: str) -> Any:
    return json.loads(text)

def parse_xml_to_data(text: str) -> dict:
    root = ET.fromstring(text)

    def recurse(node):
        item = {}
        if node.attrib:
            item["@attributes"] = dict(node.attrib)
        children = list(node)
        if children:
            grouped = {}
            for child in children:
                grouped.setdefault(child.tag, []).append(recurse(child))
            for k, v in grouped.items():
                item[k] = v if len(v) > 1 else v[0]
        txt = (node.text or "").strip()
        if txt:
            item["#text"] = txt
        return item

    return {root.tag: recurse(root)}

def data_to_xml(data: Any, root_name: str = "root") -> str:
    def build(parent, value, tag):
        if isinstance(value, dict):
            elem = ET.SubElement(parent, tag)
            for k, v in value.items():
                build(elem, v, k)
        elif isinstance(value, list):
            for item in value:
                build(parent, item, tag)
        else:
            elem = ET.SubElement(parent, tag)
            elem.text = "" if value is None else str(value)

    root = ET.Element(root_name)
    if isinstance(data, dict):
        for k, v in data.items():
            build(root, v, k)
    elif isinstance(data, list):
        for item in data:
            build(root, item, "item")
    else:
        build(root, data, "item")
    return ET.tostring(root, encoding="unicode")

def data_to_yaml(data: Any) -> str:
    return yaml.safe_dump(data, sort_keys=False, allow_unicode=True)

def data_to_json(data: Any) -> str:
    return json.dumps(data, indent=2, ensure_ascii=False)

def data_to_text(data: Any) -> str:
    if isinstance(data, str):
        return data
    if isinstance(data, list):
        return "\n".join(str(x) for x in data)
    if isinstance(data, dict):
        return json.dumps(data, indent=2, ensure_ascii=False)
    return str(data)

def split_long_text(text: str, max_chars: int = 2600) -> list[str]:
    blocks = []
    current = ""
    for para in text.splitlines() or [text]:
        if len(current) + len(para) + 1 > max_chars and current:
            blocks.append(current.strip())
            current = ""
        current += para + "\n"
    if current.strip():
        blocks.append(current.strip())
    return blocks or [""]

def make_text_image(text: str, out_png: Path, width: int = 1600, padding: int = 64, bg=(255,255,255), fg=(20,20,20)) -> Path:
    font_paths = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]
    font = None
    for fp in font_paths:
        if Path(fp).exists():
            font = ImageFont.truetype(fp, 28)
            break
    if font is None:
        font = ImageFont.load_default()
    draw = ImageDraw.Draw(Image.new("RGB", (width, 100), bg))
    lines = []
    for para in text.splitlines() or [text]:
        wrapped = textwrap.wrap(para, width=72) or [""]
        lines.extend(wrapped)
    line_h = int((font.size if hasattr(font, "size") else 22) * 1.6)
    height = max(240, padding * 2 + line_h * (len(lines) + 1))
    img = Image.new("RGB", (width, height), bg)
    draw = ImageDraw.Draw(img)
    y = padding
    for line in lines:
        draw.text((padding, y), line, fill=fg, font=font)
        y += line_h
    ensure_parent(out_png)
    img.save(out_png, optimize=True)
    return out_png

def pdf_from_content(title: str, text: str | None = None, df: pd.DataFrame | None = None, images: list[Path] | None = None, out_pdf: Path | None = None) -> Path:
    out_pdf = out_pdf or Path(tempfile.mkstemp(suffix=".pdf")[1])
    ensure_parent(out_pdf)
    doc = SimpleDocTemplate(str(out_pdf), pagesize=portrait(A4), leftMargin=36, rightMargin=36, topMargin=42, bottomMargin=36)
    styles = getSampleStyleSheet()
    story = []
    title_style = ParagraphStyle(
        "title",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=20,
        leading=24,
        alignment=TA_LEFT,
        spaceAfter=12,
    )
    body_style = ParagraphStyle(
        "body",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        spaceAfter=6,
    )
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 8))
    if text:
        for block in split_long_text(text, 1200):
            for para in block.split("\n"):
                if para.strip():
                    story.append(Paragraph(para.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), body_style))
                    story.append(Spacer(1, 4))
    if df is not None and not df.empty:
        table_df = df.head(200)  # safety for huge files
        data = [list(table_df.columns)] + table_df.astype(str).values.tolist()
        col_count = len(table_df.columns)
        col_widths = [max(40, 500 // max(col_count, 1))] * col_count
        table = Table(data, colWidths=col_widths, repeatRows=1)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#1f2937")),
            ("TEXTCOLOR", (0,0), (-1,0), colors.white),
            ("GRID", (0,0), (-1,-1), 0.25, colors.HexColor("#9ca3af")),
            ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
            ("FONTNAME", (0,1), (-1,-1), "Helvetica"),
            ("FONTSIZE", (0,0), (-1,-1), 8),
            ("VALIGN", (0,0), (-1,-1), "TOP"),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.whitesmoke, colors.HexColor("#f3f4f6")]),
        ]))
        story.append(Spacer(1, 10))
        story.append(table)
    if images:
        story.append(PageBreak())
        for img_path in images:
            try:
                img = Image.open(img_path)
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                max_w = 500
                max_h = 650
                ratio = min(max_w / img.width, max_h / img.height)
                w = img.width * ratio
                h = img.height * ratio
                story.append(Spacer(1, 12))
                story.append(Paragraph(img_path.name, body_style))
                story.append(Spacer(1, 6))
                story.append(reportlab_image(img_path, w, h))
            except Exception as exc:
                story.append(Paragraph(f"Image preview unavailable for {img_path.name}: {exc}", body_style))
    doc.build(story)
    return out_pdf

def reportlab_image(img_path: Path, w: float, h: float):
    from reportlab.platypus import Image as RLImage
    return RLImage(str(img_path), width=w, height=h)

def write_docx(title: str, text: str | None = None, df: pd.DataFrame | None = None, images: list[Path] | None = None, out_path: Path | None = None) -> Path:
    out_path = out_path or Path(tempfile.mkstemp(suffix=".docx")[1])
    ensure_parent(out_path)
    doc = Document()
    doc.add_heading(title, level=1)
    if text:
        for block in split_long_text(text, 2000):
            for para in block.split("\n"):
                if para.strip():
                    doc.add_paragraph(para)
    if df is not None and not df.empty:
        table_df = df.head(200)
        table = doc.add_table(rows=1, cols=len(table_df.columns))
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        for idx, col in enumerate(table_df.columns):
            hdr[idx].text = str(col)
        for _, row in table_df.iterrows():
            cells = table.add_row().cells
            for idx, val in enumerate(row.tolist()):
                cells[idx].text = "" if pd.isna(val) else str(val)
    if images:
        for img_path in images:
            try:
                doc.add_picture(str(img_path), width=Inches(5.8))
            except Exception:
                doc.add_paragraph(f"[Image could not be embedded: {img_path.name}]")
    doc.save(str(out_path))
    return out_path

def write_pptx(title: str, text: str | None = None, df: pd.DataFrame | None = None, images: list[Path] | None = None, out_path: Path | None = None) -> Path:
    out_path = out_path or Path(tempfile.mkstemp(suffix=".pptx")[1])
    ensure_parent(out_path)
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    # Title slide
    slide = prs.slides.add_slide(slide_layout)
    tx_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.5), PptxInches(8.5), PptxInches(0.8))
    tf = tx_box.text_frame
    tf.text = title
    if text:
        chunks = split_long_text(text, 1500)
        for chunk in chunks[:8]:
            slide = prs.slides.add_slide(slide_layout)
            box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.4), PptxInches(8.8), PptxInches(6.8))
            tf = box.text_frame
            for i, para in enumerate(chunk.split("\n")):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para
                p.font.size = Pt(18 if i == 0 else 14)
    if df is not None and not df.empty:
        slide = prs.slides.add_slide(slide_layout)
        box = slide.shapes.add_textbox(PptxInches(0.3), PptxInches(0.2), PptxInches(9.0), PptxInches(0.5))
        box.text_frame.text = f"{title} - Data Preview"
        rows = min(len(df), 10)
        cols = len(df.columns)
        table = slide.shapes.add_table(rows + 1, cols, PptxInches(0.3), PptxInches(0.9), PptxInches(9.0), PptxInches(4.8)).table
        for c, col in enumerate(df.columns):
            table.cell(0, c).text = str(col)
        for r in range(rows):
            for c, val in enumerate(df.iloc[r].tolist()):
                table.cell(r + 1, c).text = "" if pd.isna(val) else str(val)
    if images:
        for img_path in images:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(str(img_path), PptxInches(0.5), PptxInches(0.8), width=PptxInches(8.6))
    prs.save(str(out_path))
    return out_path

def write_spreadsheet(title: str, data: Any, out_path: Path | None = None, ext: str = "xlsx") -> Path:
    out_path = out_path or Path(tempfile.mkstemp(suffix=f".{ext}")[1])
    ensure_parent(out_path)
    df = dataframe_from_any(data)
    if ext == "xlsx":
        with pd.ExcelWriter(out_path, engine="openpyxl") as writer:
            df.to_excel(writer, index=False, sheet_name="Sheet1")
    elif ext == "ods":
        with pd.ExcelWriter(out_path, engine="odf") as writer:
            df.to_excel(writer, index=False, sheet_name="Sheet1")
    elif ext == "csv":
        df.to_csv(out_path, index=False)
    elif ext == "json":
        out_path.write_text(json.dumps(df.to_dict(orient="records"), indent=2, ensure_ascii=False), encoding="utf-8")
    else:
        raise RuntimeError(f"Unsupported spreadsheet output: {ext}")
    return out_path

def write_textlike(data: Any, out_path: Path, ext: str) -> Path:
    ensure_parent(out_path)
    ext = normalise_ext(ext)
    if ext == "txt":
        out_path.write_text(data_to_text(data), encoding="utf-8")
    elif ext == "json":
        out_path.write_text(data_to_json(data), encoding="utf-8")
    elif ext == "yaml":
        out_path.write_text(data_to_yaml(data), encoding="utf-8")
    elif ext == "xml":
        out_path.write_text(data_to_xml(data), encoding="utf-8")
    elif ext == "csv":
        df = dataframe_from_any(data)
        df.to_csv(out_path, index=False)
    else:
        raise RuntimeError(f"Unsupported text output: {ext}")
    return out_path

def extract_pdf_text_and_tables(path: Path) -> tuple[str, pd.DataFrame | None]:
    reader = PdfReader(str(path))
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            texts.append("")
    text = "\n\n".join([t.strip() for t in texts if t.strip()])
    df = pd.DataFrame({"page": list(range(1, len(texts)+1)), "text": texts}) if texts else None
    return text or path.stem, df

def extract_docx_text_and_tables(path: Path) -> tuple[str, pd.DataFrame | None, list[Image.Image] | None]:
    doc = Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    tables = []
    for t in doc.tables:
        rows = []
        for row in t.rows:
            rows.append([c.text.strip() for c in row.cells])
        if rows:
            tables.append(rows)
    text = "\n".join(paras)
    table_df = None
    if tables:
        # flatten first table only for simplicity
        first = tables[0]
        if len(first) >= 1:
            header = first[0]
            body = first[1:] if len(first) > 1 else []
            if body:
                table_df = pd.DataFrame(body, columns=header)
    return text or path.stem, table_df, None

def extract_pptx_text_and_tables(path: Path) -> tuple[str, pd.DataFrame | None]:
    prs = Presentation(str(path))
    lines = []
    slide_rows = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_texts.append(shape.text.strip())
        joined = "\n".join(slide_texts)
        slide_rows.append({"slide": i, "text": joined})
        if joined:
            lines.append(f"Slide {i}\n{joined}")
    df = pd.DataFrame(slide_rows) if slide_rows else None
    return "\n\n".join(lines) or path.stem, df

def extract_text_data(path: Path, ext: str) -> tuple[Any, pd.DataFrame | None, str]:
    ext = normalise_ext(ext)
    if ext == "txt":
        text = safe_utf8_text(path)
        return text, pd.DataFrame({"line": text.splitlines()}), "text"
    if ext == "csv":
        try:
            df = pd.read_csv(path)
        except Exception:
            df = pd.read_csv(path, header=None)
        return df.to_dict(orient="records"), df, "table"
    if ext == "json":
        obj = json.loads(safe_utf8_text(path))
        df = dataframe_from_any(obj)
        return obj, df, "data"
    if ext in {"yaml", "yml"}:
        obj = yaml.safe_load(safe_utf8_text(path))
        df = dataframe_from_any(obj)
        return obj, df, "data"
    if ext == "xml":
        obj = parse_xml_to_data(safe_utf8_text(path))
        df = dataframe_from_any(obj)
        return obj, df, "data"
    raise RuntimeError(f"Unsupported text/data source: {ext}")

def extract_spreadsheet(path: Path, ext: str) -> tuple[Any, pd.DataFrame, str]:
    ext = normalise_ext(ext)
    if ext == "csv":
        df = pd.read_csv(path)
        return df.to_dict(orient="records"), df, "table"
    if ext == "xlsx":
        try:
            xls = pd.read_excel(path, sheet_name=None, engine="openpyxl")
        except Exception:
            xls = pd.read_excel(path, sheet_name=None)
        # combine first sheet
        first_sheet = next(iter(xls.values()))
        df = first_sheet
        payload = {sheet: frame.fillna("").to_dict(orient="records") for sheet, frame in xls.items()}
        return payload, df, "table"
    if ext == "ods":
        xls = pd.read_excel(path, sheet_name=None, engine="odf")
        first_sheet = next(iter(xls.values()))
        df = first_sheet
        payload = {sheet: frame.fillna("").to_dict(orient="records") for sheet, frame in xls.items()}
        return payload, df, "table"
    if ext == "json":
        obj = json.loads(safe_utf8_text(path))
        df = dataframe_from_any(obj)
        return obj, df, "data"
    raise RuntimeError(f"Unsupported spreadsheet source: {ext}")

def _plain_text_from_rtf(raw_text: str) -> str:
    """Best-effort plain-text extraction from RTF markup without a full parser.

    RTF is ASCII-safe text, so it always decodes fine, but showing the raw
    control words (\\rtf1\\ansi\\deff0 ...) to a user isn't a real conversion.
    This strips the common control-word/escape syntax to leave readable text.
    """
    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", raw_text)  # hex-escaped chars
    text = re.sub(r"\\u-?\d+\??", "", text)  # unicode escapes
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)  # control words
    text = text.replace("{", "").replace("}", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_document(path: Path, ext: str) -> tuple[Any, pd.DataFrame | None, str]:
    ext = normalise_ext(ext)
    if ext == "pdf":
        text, df = extract_pdf_text_and_tables(path)
        return text, df, "document"
    if ext == "docx":
        text, df, _ = extract_docx_text_and_tables(path)
        return text, df, "document"
    if ext in {"odt", "rtf", "doc"} and HAS_LIBREOFFICE:
        # Convert to txt first, then read.
        with tempfile.TemporaryDirectory() as td:
            td = Path(td)
            try:
                subprocess.run([
                    "libreoffice", "--headless", "--convert-to", "txt:Text",
                    "--outdir", str(td), str(path)
                ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                txt_files = list(td.glob("*.txt"))
                if txt_files:
                    text = safe_utf8_text(txt_files[0])
                    return text, pd.DataFrame({"line": text.splitlines()}), "document"
            except Exception as exc:
                logger.warning("LibreOffice txt conversion failed for %s: %s", path, exc)
    if ext == "rtf":
        # RTF is plain-text markup, so a readable fallback is possible even
        # without LibreOffice.
        text = _plain_text_from_rtf(safe_utf8_text(path))
        return text, pd.DataFrame({"line": text.splitlines()}), "document"
    if ext in {"odt", "doc"}:
        # These are genuinely binary/compressed containers (doc = OLE2,
        # odt = zipped XML). Without LibreOffice there's no reliable way to
        # read them, so fail clearly instead of returning decoded garbage.
        raise RuntimeError(
            f"Converting .{ext} files requires LibreOffice, which isn't available in this "
            "environment. Try .docx, .pdf, or .rtf instead."
        )
    text = safe_utf8_text(path)
    return text, pd.DataFrame({"line": text.splitlines()}), "document"

def extract_presentation(path: Path, ext: str) -> tuple[Any, pd.DataFrame | None, str]:
    ext = normalise_ext(ext)
    if ext == "pptx":
        text, df = extract_pptx_text_and_tables(path)
        return text, df, "presentation"
    if ext in {"ppt", "odp"} and HAS_LIBREOFFICE:
        with tempfile.TemporaryDirectory() as td:
            td = Path(td)
            try:
                subprocess.run([
                    "libreoffice", "--headless", "--convert-to", "pptx",
                    "--outdir", str(td), str(path)
                ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                pptx_files = list(td.glob("*.pptx"))
                if pptx_files:
                    text, df = extract_pptx_text_and_tables(pptx_files[0])
                    return text, df, "presentation"
            except Exception as exc:
                logger.warning("LibreOffice pptx conversion failed for %s: %s", path, exc)
    if ext in {"ppt", "odp"}:
        # Both are binary/compressed containers (ppt = OLE2, odp = zipped
        # XML) — without LibreOffice there's no reliable way to read them.
        raise RuntimeError(
            f"Converting .{ext} files requires LibreOffice, which isn't available in this "
            "environment. Try .pptx instead."
        )
    text = safe_utf8_text(path)
    return text, pd.DataFrame({"line": text.splitlines()}), "presentation"

def extract_image_payload(path: Path, ext: str) -> tuple[Any, pd.DataFrame | None, str, list[Path]]:
    img_path = path
    if normalise_ext(ext) == "svg":
        tmp_png = path.with_suffix(".render.png")
        image_from_svg(path, tmp_png)
        img_path = tmp_png
    img = load_image_any(img_path)
    return {"image": img}, None, "image", [img_path]

def build_capabilities() -> dict[str, list[str]]:
    # Route every source through the same environment-aware filter used at
    # conversion time, so this never advertises a target that will 500.
    return {ext: capabilities_for_source(ext) for ext in SUPPORTED_CONVERSIONS}

# -----------------------------
# Output creators
# -----------------------------
def create_image_output(payload: dict[str, Any], target_ext: str, out_path: Path) -> Path:
    target_ext = normalise_ext(target_ext)
    ensure_parent(out_path)
    if payload.get("images"):
        img_path = Path(payload["images"][0])
        img = load_image_any(img_path)
    else:
        # fallback: render text/data to a tall image
        text = payload.get("text") or json.dumps(payload.get("data") or {}, indent=2, ensure_ascii=False)
        return make_text_image(text, out_path) if target_ext == "png" else image_to_raster(Image.open(make_text_image(text, out_path.with_suffix(".png"))), out_path, target_ext)

    if target_ext in {"png", "jpg", "jpeg", "webp", "gif", "bmp", "tif", "tiff"}:
        return image_to_raster(img, out_path, target_ext)
    if target_ext == "pdf":
        return image_to_pdf(img, out_path)
    if target_ext == "docx":
        return write_docx(payload["title"], text=payload.get("text"), images=[img_path], out_path=out_path)
    if target_ext == "pptx":
        return write_pptx(payload["title"], text=payload.get("text"), images=[img_path], out_path=out_path)
    raise RuntimeError(f"Unsupported image output format: {target_ext}")

def create_text_output(payload: dict[str, Any], target_ext: str, out_path: Path) -> Path:
    target_ext = normalise_ext(target_ext)
    data = payload.get("data")
    df = payload.get("df")
    if target_ext in {"txt", "json", "yaml", "xml", "csv"}:
        return write_textlike(data if data is not None else payload.get("text", ""), out_path, target_ext)
    if target_ext in {"xlsx", "ods"}:
        return write_spreadsheet(payload["title"], data if data is not None else payload.get("text", ""), out_path, target_ext)
    if target_ext == "pdf":
        return pdf_from_content(payload["title"], text=payload.get("text"), df=df, out_pdf=out_path)
    if target_ext == "docx":
        return write_docx(payload["title"], text=payload.get("text"), df=df, out_path=out_path)
    if target_ext == "pptx":
        return write_pptx(payload["title"], text=payload.get("text"), df=df, out_path=out_path)
    if target_ext == "png":
        return make_text_image(payload.get("text") or data_to_text(data), out_path)
    raise RuntimeError(f"Unsupported text/data output format: {target_ext}")

def create_spreadsheet_output(payload: dict[str, Any], target_ext: str, out_path: Path) -> Path:
    target_ext = normalise_ext(target_ext)
    data = payload.get("data")
    df = payload.get("df")
    if df is None:
        df = dataframe_from_any(data)
    if target_ext in {"xlsx", "ods", "csv", "json"}:
        return write_spreadsheet(payload["title"], df, out_path, target_ext)
    if target_ext == "txt":
        return write_textlike(df.to_dict(orient="records"), out_path, "txt")
    if target_ext == "pdf":
        return pdf_from_content(payload["title"], text=payload.get("text"), df=df, out_pdf=out_path)
    if target_ext == "docx":
        return write_docx(payload["title"], text=payload.get("text"), df=df, out_path=out_path)
    if target_ext == "pptx":
        return write_pptx(payload["title"], text=payload.get("text"), df=df, out_path=out_path)
    if target_ext == "png":
        return make_text_image(df.to_string(index=False), out_path)
    raise RuntimeError(f"Unsupported spreadsheet output format: {target_ext}")

def create_document_output(payload: dict[str, Any], target_ext: str, out_path: Path) -> Path:
    target_ext = normalise_ext(target_ext)
    text = payload.get("text") or payload["title"]
    df = payload.get("df")
    if target_ext == "pdf":
        return pdf_from_content(payload["title"], text=text, df=df, out_pdf=out_path)
    if target_ext == "docx":
        return write_docx(payload["title"], text=text, df=df, out_path=out_path)
    if target_ext == "txt":
        return write_textlike(text, out_path, "txt")
    if target_ext == "csv":
        return write_spreadsheet(payload["title"], df if df is not None else pd.DataFrame({"line": text.splitlines()}), out_path, "csv")
    if target_ext == "json":
        return write_textlike({"title": payload["title"], "text": text, "rows": df.to_dict(orient="records") if df is not None else []}, out_path, "json")
    if target_ext == "xlsx":
        return write_spreadsheet(payload["title"], df if df is not None else pd.DataFrame({"line": text.splitlines()}), out_path, "xlsx")
    if target_ext == "pptx":
        return write_pptx(payload["title"], text=text, df=df, out_path=out_path)
    if target_ext == "png":
        # create a PDF first, then render first page to PNG
        tmp_pdf = out_path.with_suffix(".tmp.pdf")
        pdf_from_content(payload["title"], text=text, df=df, out_pdf=tmp_pdf)
        images = render_pdf_to_images(tmp_pdf, out_path.parent, out_path.stem)
        if len(images) == 1:
            images[0].replace(out_path)
            tmp_pdf.unlink(missing_ok=True)
            return out_path
        zip_paths(images, out_path.with_suffix(".zip"))
        tmp_pdf.unlink(missing_ok=True)
        return out_path.with_suffix(".zip")
    raise RuntimeError(f"Unsupported document output format: {target_ext}")

def create_presentation_output(payload: dict[str, Any], target_ext: str, out_path: Path) -> Path:
    target_ext = normalise_ext(target_ext)
    text = payload.get("text") or payload["title"]
    df = payload.get("df")
    if target_ext == "pptx":
        return write_pptx(payload["title"], text=text, df=df, out_path=out_path)
    if target_ext == "pdf":
        return pdf_from_content(payload["title"], text=text, df=df, out_pdf=out_path)
    if target_ext == "docx":
        return write_docx(payload["title"], text=text, df=df, out_path=out_path)
    if target_ext == "txt":
        return write_textlike(text, out_path, "txt")
    if target_ext == "csv":
        return write_spreadsheet(payload["title"], df if df is not None else pd.DataFrame({"slide": text.splitlines()}), out_path, "csv")
    if target_ext == "json":
        return write_textlike({"title": payload["title"], "slides": text.splitlines()}, out_path, "json")
    if target_ext == "xlsx":
        return write_spreadsheet(payload["title"], df if df is not None else pd.DataFrame({"slide": text.splitlines()}), out_path, "xlsx")
    if target_ext == "png":
        tmp_pdf = out_path.with_suffix(".tmp.pdf")
        pdf_from_content(payload["title"], text=text, df=df, out_pdf=tmp_pdf)
        images = render_pdf_to_images(tmp_pdf, out_path.parent, out_path.stem)
        if len(images) == 1:
            images[0].replace(out_path)
            tmp_pdf.unlink(missing_ok=True)
            return out_path
        zip_paths(images, out_path.with_suffix(".zip"))
        tmp_pdf.unlink(missing_ok=True)
        return out_path.with_suffix(".zip")
    raise RuntimeError(f"Unsupported presentation output format: {target_ext}")

def create_media_output(source_path: Path, source_ext: str, target_ext: str, out_path: Path, black_screen: bool = True) -> Path:
    if not HAS_FFMPEG:
        raise RuntimeError("ffmpeg is required for media conversions")
    ensure_parent(out_path)
    source_ext = normalise_ext(source_ext)
    target_ext = normalise_ext(target_ext)
    if source_ext in VIDEO_EXTS and target_ext in AUDIO_EXTS:
        cmd = [
            "ffmpeg", "-y", "-i", str(source_path),
            "-vn", "-acodec", "libmp3lame" if target_ext == "mp3" else ("pcm_s16le" if target_ext == "wav" else "libvorbis"),
        ]
        if target_ext == "mp3":
            cmd += ["-q:a", "2"]
        cmd += [str(out_path)]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return out_path
    if source_ext in AUDIO_EXTS and target_ext in VIDEO_EXTS:
        screen = "black" if black_screen else "white"
        cmd = [
            "ffmpeg", "-y",
            "-f", "lavfi", "-i", f"color=c={screen}:s=1280x720:r=30",
            "-i", str(source_path),
            "-shortest",
            "-c:v", "libx264", "-pix_fmt", "yuv420p",
            "-c:a", "aac",
            str(out_path),
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return out_path
    if source_ext in VIDEO_EXTS and target_ext in VIDEO_EXTS:
        cmd = ["ffmpeg", "-y", "-i", str(source_path)]
        if target_ext == "mp4":
            cmd += ["-c:v", "libx264", "-c:a", "aac", str(out_path)]
        elif target_ext == "webm":
            cmd += ["-c:v", "libvpx-vp9", "-c:a", "libopus", str(out_path)]
        elif target_ext == "mov":
            cmd += ["-c:v", "prores_ks", "-profile:v", "3", str(out_path)]
        elif target_ext == "avi":
            cmd += ["-c:v", "mpeg4", "-c:a", "mp3", str(out_path)]
        elif target_ext == "mkv":
            cmd += ["-c:v", "libx264", "-c:a", "aac", str(out_path)]
        else:
            raise RuntimeError(f"Unsupported video target: {target_ext}")
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return out_path
    if source_ext in AUDIO_EXTS and target_ext in AUDIO_EXTS:
        if source_ext == target_ext:
            shutil.copy2(source_path, out_path)
            return out_path
        codec_map = {
            "mp3": "libmp3lame",
            "wav": "pcm_s16le",
            "ogg": "libvorbis",
            "flac": "flac",
        }
        cmd = ["ffmpeg", "-y", "-i", str(source_path), "-c:a", codec_map[target_ext], str(out_path)]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return out_path
    raise RuntimeError(f"Unsupported media conversion: {source_ext} -> {target_ext}")

def _ensure_safe_archive_member(name: str, extract_dir: Path) -> None:
    """Reject archive entries that would write outside extract_dir ("Zip Slip")."""
    if not name or name.startswith(("/", "\\")):
        raise RuntimeError("Archive contains an unsafe absolute path and was rejected.")
    dest = (extract_dir / name).resolve()
    if dest != extract_dir and extract_dir not in dest.parents:
        raise RuntimeError("Archive contains a path that escapes the extraction folder and was rejected.")


def extract_archive(archive_path: Path, out_zip: Path) -> Path:
    ensure_parent(out_zip)
    extract_dir = out_zip.parent / f"{out_zip.stem}_extract"
    extract_dir.mkdir(parents=True, exist_ok=True)
    extract_dir_resolved = extract_dir.resolve()

    if zipfile.is_zipfile(archive_path):
        with zipfile.ZipFile(archive_path) as zf:
            for name in zf.namelist():
                _ensure_safe_archive_member(name, extract_dir_resolved)
            zf.extractall(extract_dir)
    elif tarfile.is_tarfile(archive_path):
        with tarfile.open(archive_path) as tf:
            members = [m for m in tf.getmembers() if m.isfile() or m.isdir()]
            for member in members:
                _ensure_safe_archive_member(member.name, extract_dir_resolved)
            tf.extractall(extract_dir, members=members)  # noqa: S202 - members vetted above
    else:
        raise RuntimeError("Only ZIP and TAR archives are supported in this build.")
    # Bundle extracted files as a zip for the browser
    extracted_files = [p for p in extract_dir.rglob("*") if p.is_file()]
    if not extracted_files:
        raise RuntimeError("Archive is empty or could not be extracted.")
    with zipfile.ZipFile(out_zip, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for file in extracted_files:
            zf.write(file, arcname=str(file.relative_to(extract_dir)))
    shutil.rmtree(extract_dir, ignore_errors=True)
    return out_zip

def capabilities_for_source(source_ext: str) -> list[str]:
    source_ext = normalise_ext(source_ext)
    supported = SUPPORTED_CONVERSIONS.get(source_ext, []).copy()
    # Keep the advertised list honest about what this environment can actually do.
    if not HAS_FFMPEG and (source_ext in VIDEO_EXTS or source_ext in AUDIO_EXTS):
        # create_media_output() requires ffmpeg unconditionally, so no video/audio
        # target (not even same-category, e.g. mp4->webm) will work without it.
        supported = [ext for ext in supported if ext not in VIDEO_EXTS and ext not in AUDIO_EXTS]
    if not HAS_INKSCAPE and not HAS_MAGICK and source_ext == "svg":
        # Every SVG target (including pdf/docx/pptx) is rasterized via
        # image_from_svg() first, so with neither tool, nothing works.
        supported = []
    if not HAS_LIBREOFFICE and source_ext in {"doc", "odt", "ppt", "odp"}:
        # These are binary/compressed containers that extract_document() /
        # extract_presentation() can only read via LibreOffice.
        supported = []
    return supported

def detect_output_mime(ext: str) -> str:
    ext = normalise_ext(ext)
    return MIME_BY_EXT.get(ext, mimetypes.guess_type(f"file.{ext}")[0] or "application/octet-stream")

def write_output_bundle(name_base: str, files: list[Path], out_zip: Path) -> Path:
    ensure_parent(out_zip)
    with zipfile.ZipFile(out_zip, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for file in files:
            zf.write(file, arcname=file.name)
    return out_zip

def make_payload_from_source(source_path: Path, source_ext: str) -> dict[str, Any]:
    source_ext = normalise_ext(source_ext)
    kind = source_category(source_ext)
    title = stem_of(source_path.name)
    if kind == "image":
        if source_ext == "svg":
            tmp_png = source_path.with_suffix(".render.png")
            image_from_svg(source_path, tmp_png)
            image_paths = [tmp_png]
        else:
            image_paths = [source_path]
        img = load_image_any(image_paths[0])
        return {
            "kind": kind,
            "title": title,
            "text": f"Image conversion for {source_path.name}",
            "data": {"file": source_path.name, "type": source_ext},
            "df": None,
            "images": image_paths,
            "image": img,
        }
    if kind == "text-data":
        data, df, data_kind = extract_text_data(source_path, source_ext)
        return {
            "kind": kind,
            "title": title,
            "text": data_to_text(data),
            "data": data,
            "df": df,
            "images": [],
        }
    if kind == "spreadsheet":
        data, df, data_kind = extract_spreadsheet(source_path, source_ext)
        return {
            "kind": kind,
            "title": title,
            "text": data_to_text(data),
            "data": data,
            "df": df,
            "images": [],
        }
    if kind == "document":
        data, df, _ = extract_document(source_path, source_ext)
        return {
            "kind": kind,
            "title": title,
            "text": data_to_text(data),
            "data": data,
            "df": df,
            "images": [],
        }
    if kind == "presentation":
        data, df, _ = extract_presentation(source_path, source_ext)
        return {
            "kind": kind,
            "title": title,
            "text": data_to_text(data),
            "data": data,
            "df": df,
            "images": [],
        }
    if kind in {"video", "audio"}:
        return {
            "kind": kind,
            "title": title,
            "text": f"{kind.title()} file: {source_path.name}",
            "data": {"file": source_path.name, "type": source_ext},
            "df": None,
            "images": [],
        }
    if kind == "archive":
        return {
            "kind": kind,
            "title": title,
            "text": f"Archive file: {source_path.name}",
            "data": {"file": source_path.name, "type": source_ext},
            "df": None,
            "images": [],
        }
    raise RuntimeError(f"Unsupported source file type: {source_ext}")

def convert_file(source_path: Path, source_ext: str, target_ext: str, session_dirs: dict[str, Path]) -> dict[str, Any]:
    source_ext = normalise_ext(source_ext)
    target_ext = normalise_ext(target_ext)
    payload = make_payload_from_source(source_path, source_ext)
    kind = payload["kind"]
    base_stem = stem_of(source_path.name)
    safe_target = "extract" if target_ext == "extract" else target_ext
    output_dir = session_dirs["outputs"]

    # Output file naming
    if kind == "archive" and safe_target == "extract":
        out_zip = output_dir / make_output_name(base_stem, "extracted", "zip")
        final_path = extract_archive(source_path, out_zip)
        return {
            "output_path": final_path,
            "output_file_name": final_path.name,
            "download_name": final_path.name,
            "output_mime": "application/zip",
            "output_size": final_path.stat().st_size,
        }

    # media
    if kind in {"video", "audio"}:
        out_path = output_dir / make_output_name(base_stem, "converted", safe_target)
        final_path = create_media_output(source_path, source_ext, target_ext, out_path)
        return {
            "output_path": final_path,
            "output_file_name": final_path.name,
            "download_name": final_path.name,
            "output_mime": detect_output_mime(safe_target),
            "output_size": final_path.stat().st_size,
        }

    # image
    if kind == "image":
        out_path = output_dir / make_output_name(base_stem, "converted", safe_target)
        final_path = create_image_output(payload, safe_target, out_path)
        return {
            "output_path": final_path,
            "output_file_name": final_path.name,
            "download_name": final_path.name,
            "output_mime": detect_output_mime(safe_target if final_path.suffix != ".zip" else "zip"),
            "output_size": final_path.stat().st_size,
        }

    # text/data
    if kind == "text-data":
        out_path = output_dir / make_output_name(base_stem, "converted", safe_target)
        final_path = create_text_output(payload, safe_target, out_path)
        return {
            "output_path": final_path,
            "output_file_name": final_path.name,
            "download_name": final_path.name,
            "output_mime": detect_output_mime(final_path.suffix.lstrip(".")),
            "output_size": final_path.stat().st_size,
        }

    # spreadsheet
    if kind == "spreadsheet":
        out_path = output_dir / make_output_name(base_stem, "converted", safe_target)
        final_path = create_spreadsheet_output(payload, safe_target, out_path)
        return {
            "output_path": final_path,
            "output_file_name": final_path.name,
            "download_name": final_path.name,
            "output_mime": detect_output_mime(final_path.suffix.lstrip(".")),
            "output_size": final_path.stat().st_size,
        }

    # document
    if kind == "document":
        out_path = output_dir / make_output_name(base_stem, "converted", safe_target)
        final_path = create_document_output(payload, safe_target, out_path)
        return {
            "output_path": final_path,
            "output_file_name": final_path.name,
            "download_name": final_path.name,
            "output_mime": detect_output_mime(final_path.suffix.lstrip(".")),
            "output_size": final_path.stat().st_size,
        }

    # presentation
    if kind == "presentation":
        out_path = output_dir / make_output_name(base_stem, "converted", safe_target)
        final_path = create_presentation_output(payload, safe_target, out_path)
        return {
            "output_path": final_path,
            "output_file_name": final_path.name,
            "download_name": final_path.name,
            "output_mime": detect_output_mime(final_path.suffix.lstrip(".")),
            "output_size": final_path.stat().st_size,
        }

    raise RuntimeError(f"Unsupported conversion path: {source_ext} -> {target_ext}")
